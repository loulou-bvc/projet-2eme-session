#!/usr/bin/env python3
"""
Génération de la présentation PowerPoint - Rôle 1 : Data Engineering
Projet 8 : Prix Négatifs de l'Électricité Renouvelable en Europe
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ============================================================
# COULEURS & STYLES
# ============================================================
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)        # Fond sombre principal
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xC7)     # Bleu accent
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)    # Vert accent
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)   # Orange accent
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)      # Rouge accent
ACCENT_YELLOW = RGBColor(0xF1, 0xC4, 0x0F)   # Jaune
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBD, 0xBD, 0xBD)
MEDIUM_GRAY = RGBColor(0x75, 0x75, 0x75)
DARK_CARD = RGBColor(0x25, 0x25, 0x40)       # Fond carte
VERY_DARK = RGBColor(0x12, 0x12, 0x22)       # Fond très sombre
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
SOFT_BG = RGBColor(0xF5, 0xF5, 0xF5)         # Fond clair alternatif
NAVY = RGBColor(0x0D, 0x47, 0xA1)            # Bleu marine
TEAL = RGBColor(0x00, 0x89, 0x7B)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)


def set_slide_bg(slide, color):
    """Définir la couleur de fond d'une diapositive."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, transparency=0):
    """Ajouter un rectangle coloré."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    """Ajouter un rectangle arrondi."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Ajouter une zone de texte."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(3),
                  font_name="Calibri"):
    """Ajouter un paragraphe à un text_frame existant."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=WHITE, title=None, title_color=ACCENT_BLUE, title_size=18):
    """Ajouter une liste à puces avec titre optionnel."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.space_after = Pt(8)

        for item in items:
            add_paragraph(tf, f"  {item}", font_size, color, space_before=Pt(4), space_after=Pt(2))
    else:
        p = tf.paragraphs[0]
        p.text = f"  {items[0]}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        for item in items[1:]:
            add_paragraph(tf, f"  {item}", font_size, color, space_before=Pt(4), space_after=Pt(2))

    return txBox


def create_metric_card(slide, left, top, width, height, value, label,
                       value_color=ACCENT_BLUE, bg_color=DARK_CARD):
    """Créer une carte de métrique (nombre + label)."""
    card = add_rounded_rect(slide, left, top, width, height, bg_color)

    add_textbox(slide, left, top + Inches(0.15), width, Inches(0.6),
                value, font_size=28, color=value_color, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, left, top + Inches(0.7), width, Inches(0.5),
                label, font_size=11, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    return card


# ============================================================
# CREATION DE LA PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ============================================================
# SLIDE 1 : PAGE DE TITRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

# Bande décorative en haut
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# Bande décorative latérale
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), ACCENT_BLUE)

# Titre principal
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
            "PROJET 8", font_size=20, color=ACCENT_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
            "Prix Negatifs de l'Electricite\nRenouvelable en Europe",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Ligne séparatrice
add_shape(slide, Inches(4.5), Inches(3.6), Inches(4.333), Inches(0.04), ACCENT_BLUE)

# Sous-titre
add_textbox(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.6),
            "Role 1 : Data Engineering & Ingestion des Donnees",
            font_size=24, color=ACCENT_ORANGE, bold=False, alignment=PP_ALIGN.CENTER)

# Infos
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Semaines 1-4  |  Fevrier 2026",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Cartes en bas
create_metric_card(slide, Inches(2.5), Inches(5.8), Inches(2), Inches(1.1),
                   "50,401", "Observations horaires", ACCENT_BLUE)
create_metric_card(slide, Inches(5.0), Inches(5.8), Inches(2), Inches(1.1),
                   "71", "Variables nettoyees", ACCENT_GREEN)
create_metric_card(slide, Inches(7.5), Inches(5.8), Inches(2), Inches(1.1),
                   "0%", "Valeurs manquantes", ACCENT_GREEN)


# ============================================================
# SLIDE 2 : SOMMAIRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "SOMMAIRE", font_size=32, color=WHITE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.04), ACCENT_BLUE)

items_left = [
    ("01", "Contexte & Problematique", ACCENT_BLUE),
    ("02", "Sources de Donnees", ACCENT_GREEN),
    ("03", "Pipeline d'Ingestion", ACCENT_ORANGE),
    ("04", "Analyse de Qualite", ACCENT_RED),
    ("05", "Nettoyage des Donnees", TEAL),
]

items_right = [
    ("06", "Resultats & Metriques", PURPLE),
    ("07", "Prix Negatifs Identifies", ACCENT_YELLOW),
    ("08", "Documentation & Livrables", ACCENT_BLUE),
    ("09", "Recommandations pour l'Equipe", ACCENT_GREEN),
    ("10", "Conclusion & Prochaines Etapes", ACCENT_ORANGE),
]

for i, (num, title, color) in enumerate(items_left):
    y = Inches(1.6 + i * 1.05)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.85), DARK_CARD)
    add_textbox(slide, Inches(1.0), y + Inches(0.15), Inches(0.8), Inches(0.55),
                num, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.9), y + Inches(0.18), Inches(4.2), Inches(0.55),
                title, font_size=18, color=WHITE)

for i, (num, title, color) in enumerate(items_right):
    y = Inches(1.6 + i * 1.05)
    card = add_rounded_rect(slide, Inches(7.0), y, Inches(5.5), Inches(0.85), DARK_CARD)
    add_textbox(slide, Inches(7.2), y + Inches(0.15), Inches(0.8), Inches(0.55),
                num, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.1), y + Inches(0.18), Inches(4.2), Inches(0.55),
                title, font_size=18, color=WHITE)


# ============================================================
# SLIDE 3 : CONTEXTE & PROBLEMATIQUE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "01", font_size=14, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Contexte & Problematique", font_size=32, color=WHITE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_BLUE)

# Colonne gauche - Le phenomene
card1 = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.2), DARK_CARD)

add_textbox(slide, Inches(1.1), Inches(1.85), Inches(5.2), Inches(0.5),
            "Le Phenomene des Prix Negatifs", font_size=22, color=ACCENT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.2), Inches(4.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Quand la production eolienne et solaire depasse la demande, les prix de l'electricite deviennent NEGATIFS."
p.font.size = Pt(15)
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.space_after = Pt(12)

add_paragraph(tf, "Les producteurs paient pour ecouler leur electricite !", 15, ACCENT_ORANGE, True, space_after=Pt(16))

add_paragraph(tf, "Pourquoi c'est important ?", 17, ACCENT_GREEN, True, space_after=Pt(8))
for item in [
    "Transition energetique : phenomene en forte croissance",
    "Impact financier : pertes pour producteurs, gains pour consommateurs",
    "Gestion du reseau : defi pour les operateurs (equilibre offre-demande)",
    "Opportunite de trading : anticiper = profit"
]:
    add_paragraph(tf, f"  > {item}", 13, LIGHT_GRAY, space_before=Pt(3), space_after=Pt(3))

# Colonne droite - Objectif
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.3), DARK_CARD)

add_textbox(slide, Inches(7.1), Inches(1.85), Inches(5.2), Inches(0.5),
            "Objectif du Projet", font_size=22, color=ACCENT_GREEN, bold=True)

txBox2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.2), Inches(1.3))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Predire QUAND les prix de l'electricite deviendront negatifs en utilisant les donnees de production renouvelable, de consommation et de marche."
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.font.name = "Calibri"

# Mon role
card3 = add_rounded_rect(slide, Inches(6.8), Inches(4.2), Inches(5.7), Inches(2.7), DARK_CARD)

add_textbox(slide, Inches(7.1), Inches(4.35), Inches(5.2), Inches(0.5),
            "Mon Role : Data Engineer", font_size=22, color=ACCENT_ORANGE, bold=True)

txBox3 = slide.shapes.add_textbox(Inches(7.1), Inches(5.0), Inches(5.2), Inches(1.7))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Responsabilites :"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for item in [
    "Recherche et selection des jeux de donnees",
    "Pipeline d'ingestion reproductible (Python)",
    "Analyse de qualite et nettoyage des donnees",
    "Documentation complete (dictionnaire + rapport)"
]:
    add_paragraph(tf3, f"  > {item}", 13, LIGHT_GRAY, space_before=Pt(3), space_after=Pt(2))


# ============================================================
# SLIDE 4 : SOURCES DE DONNEES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "02", font_size=14, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Sources de Donnees", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_GREEN)

# Source 1 - OPSD Time Series (principale)
card = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(3.7), Inches(5.2), DARK_CARD)
# Badge
badge1 = add_rounded_rect(slide, Inches(0.9), Inches(1.8), Inches(1.8), Inches(0.35), ACCENT_BLUE)
add_textbox(slide, Inches(0.9), Inches(1.8), Inches(1.8), Inches(0.35),
            "PRINCIPALE", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.0), Inches(2.3), Inches(3.3), Inches(0.5),
            "OPSD Time Series", font_size=20, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(1.0), Inches(2.75), Inches(3.3), Inches(0.4),
            "v2020-10-06", font_size=13, color=LIGHT_GRAY)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.3), Inches(3.3), Inches(3.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Open Power System Data"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.font.name = "Calibri"

for item in [
    "124 Mo de donnees CSV",
    "50,401 timestamps horaires",
    "300 colonnes / 32 pays",
    "Periode : 2015-2020",
    "Licence CC-BY 4.0",
    "",
    "Variables cles :",
    "  - Prix day-ahead (EUR/MWh)",
    "  - Charge electrique (MW)",
    "  - Generation solaire/eolienne",
    "  - Capacites installees",
    "  - Profils de production",
]:
    add_paragraph(tf, item, 12, WHITE if "Variables" not in item else ACCENT_ORANGE,
                  "Variables" in item, space_before=Pt(2), space_after=Pt(1))

# Source 2 - OPSD Weather
card2 = add_rounded_rect(slide, Inches(4.8), Inches(1.7), Inches(3.7), Inches(5.2), DARK_CARD)
badge2 = add_rounded_rect(slide, Inches(4.9), Inches(1.8), Inches(2.0), Inches(0.35), ACCENT_ORANGE)
add_textbox(slide, Inches(4.9), Inches(1.8), Inches(2.0), Inches(0.35),
            "COMPLEMENTAIRE", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(5.0), Inches(2.3), Inches(3.3), Inches(0.5),
            "OPSD Weather Data", font_size=20, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(5.0), Inches(2.75), Inches(3.3), Inches(0.4),
            "ERA5 Reanalysis", font_size=13, color=LIGHT_GRAY)

txBox2 = slide.shapes.add_textbox(Inches(5.0), Inches(3.3), Inches(3.3), Inches(3.4))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Donnees meteorologiques"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.font.name = "Calibri"

for item in [
    "Plusieurs Go de donnees",
    "Grille NUTS-2 europeenne",
    "Resolution horaire",
    "",
    "Variables :",
    "  - Temperature",
    "  - Radiation solaire",
    "  - Vitesse du vent",
    "  - Precipitations",
    "  - Densite de l'air",
    "",
    "Statut : Disponible si besoin"
]:
    add_paragraph(tf2, item, 12, WHITE if "Variables" not in item and "Statut" not in item
                  else ACCENT_ORANGE if "Variables" in item else ACCENT_BLUE,
                  "Variables" in item or "Statut" in item,
                  space_before=Pt(2), space_after=Pt(1))

# Source 3 - ENTSO-E
card3 = add_rounded_rect(slide, Inches(8.8), Inches(1.7), Inches(3.7), Inches(5.2), DARK_CARD)
badge3 = add_rounded_rect(slide, Inches(8.9), Inches(1.8), Inches(2.0), Inches(0.35), MEDIUM_GRAY)
add_textbox(slide, Inches(8.9), Inches(1.8), Inches(2.0), Inches(0.35),
            "POTENTIELLE", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(9.0), Inches(2.3), Inches(3.3), Inches(0.5),
            "ENTSO-E Platform", font_size=20, color=LIGHT_GRAY, bold=True)
add_textbox(slide, Inches(9.0), Inches(2.75), Inches(3.3), Inches(0.4),
            "Transparency Platform", font_size=13, color=LIGHT_GRAY)

txBox3 = slide.shapes.add_textbox(Inches(9.0), Inches(3.3), Inches(3.3), Inches(3.4))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Donnees de marche detaillees"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.font.name = "Calibri"

for item in [
    "Generation par type (nucleaire,",
    "  gaz, charbon, eolien, solaire...)",
    "Flux transfrontaliers",
    "Prix day-ahead + intraday",
    "Donnees d'equilibrage",
    "",
    "Acces : Inscription requise",
    "API Token necessaire",
    "",
    "Statut : Disponible si besoin",
    "pour enrichir l'analyse"
]:
    add_paragraph(tf3, item, 12, WHITE if "Statut" not in item and "Acces" not in item
                  else ACCENT_BLUE if "Statut" in item else ACCENT_ORANGE,
                  "Statut" in item or "Acces" in item,
                  space_before=Pt(2), space_after=Pt(1))


# ============================================================
# SLIDE 5 : PAYS FOCUS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "02b", font_size=14, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Pays Focus : Strategie de Selection", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_GREEN)

# Allemagne
card_de = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(3.7), Inches(5.0), DARK_CARD)
flag_de = add_rounded_rect(slide, Inches(0.9), Inches(1.8), Inches(3.5), Inches(0.5), ACCENT_RED)
add_textbox(slide, Inches(0.9), Inches(1.82), Inches(3.5), Inches(0.5),
            "ALLEMAGNE (DE)", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(3.2), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Premier marche europeen"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for item in [
    "",
    "484 prix negatifs (2.76%)",
    "Min : -90.01 EUR/MWh",
    "41 colonnes de donnees",
    "",
    "Plus grande capacite renouvelable",
    "d'Europe",
    "",
    "4 operateurs reseau (TSO):",
    "  50hertz, Amprion,",
    "  TenneT, TransnetBW",
    "",
    "Variable cible principale"
]:
    is_highlight = "484" in item or "Variable cible" in item
    add_paragraph(tf, item, 13,
                  ACCENT_ORANGE if is_highlight else WHITE,
                  is_highlight,
                  space_before=Pt(2), space_after=Pt(1))

# Danemark
card_dk = add_rounded_rect(slide, Inches(4.8), Inches(1.7), Inches(3.7), Inches(5.0), DARK_CARD)
flag_dk = add_rounded_rect(slide, Inches(4.9), Inches(1.8), Inches(3.5), Inches(0.5), ACCENT_BLUE)
add_textbox(slide, Inches(4.9), Inches(1.82), Inches(3.5), Inches(0.5),
            "DANEMARK (DK)", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox2 = slide.shapes.add_textbox(Inches(5.1), Inches(2.5), Inches(3.2), Inches(4.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Champion de l'eolien"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for item in [
    "",
    "DK_1 : 539 prix negatifs (1.07%)",
    "DK_2 : 354 prix negatifs (0.70%)",
    "Min : -58.80 EUR/MWh",
    "24 colonnes de donnees",
    "",
    "Plus forte penetration eolienne",
    "en Europe (>40% du mix)",
    "",
    "2 zones de marche :",
    "  DK_1 (Est) / DK_2 (Ouest)",
    "",
    "Eolien offshore dominant"
]:
    is_highlight = "539" in item or "354" in item
    add_paragraph(tf2, item, 13,
                  ACCENT_ORANGE if is_highlight else WHITE,
                  is_highlight,
                  space_before=Pt(2), space_after=Pt(1))

# France
card_fr = add_rounded_rect(slide, Inches(8.8), Inches(1.7), Inches(3.7), Inches(5.0), DARK_CARD)
flag_fr = add_rounded_rect(slide, Inches(8.9), Inches(1.8), Inches(3.5), Inches(0.5), NAVY)
add_textbox(slide, Inches(8.9), Inches(1.82), Inches(3.5), Inches(0.5),
            "FRANCE (FR)", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox3 = slide.shapes.add_textbox(Inches(9.1), Inches(2.5), Inches(3.2), Inches(4.0))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Contraste nucleaire"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for item in [
    "",
    "0 prix negatifs (0.00%)",
    "Min : 5.00 EUR/MWh",
    "5 colonnes de donnees",
    "",
    "Domine par le nucleaire",
    "(~70% du mix electrique)",
    "",
    "Marche couple IT_NORD-FR",
    "",
    "Contraste interessant :",
    "  pourquoi pas de prix",
    "  negatifs en France ?"
]:
    is_highlight = "0 prix" in item
    add_paragraph(tf3, item, 13,
                  ACCENT_GREEN if is_highlight else WHITE,
                  is_highlight,
                  space_before=Pt(2), space_after=Pt(1))


# ============================================================
# SLIDE 6 : PIPELINE D'INGESTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "03", font_size=14, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Pipeline d'Ingestion : 4 Scripts Python", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_ORANGE)

# Script 1
scripts = [
    ("01", "download_opsd_data.py", "Telechargement", "~150 lignes",
     ACCENT_BLUE,
     ["Telechargement automatise depuis OPSD",
      "Verification d'integrite du fichier",
      "Gestion des erreurs reseau",
      "Resultat : 124 Mo CSV"]),
    ("02", "initial_exploration.py", "Exploration", "~200 lignes",
     ACCENT_GREEN,
     ["Vue d'ensemble du dataset",
      "Dimensions : 50,401 x 300",
      "Detection initiale prix negatifs",
      "Rapport : initial_exploration.txt"]),
    ("03", "data_quality_analysis.py", "Qualite", "~350 lignes",
     ACCENT_ORANGE,
     ["Analyse exhaustive de qualite",
      "26.2% valeurs manquantes globales",
      "19 colonnes >50% missing identifiees",
      "Rapport JSON automatise"]),
    ("04", "data_cleaning.py", "Nettoyage", "~250 lignes",
     ACCENT_RED,
     ["Filtrage 3 pays focus (300 -> 71 cols)",
      "Forward fill + backward fill",
      "Creation 7 variables temporelles",
      "Dataset final : 0% missing"]),
]

for i, (num, filename, phase, lines, color, details) in enumerate(scripts):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.7), Inches(2.9), Inches(5.2), DARK_CARD)

    # Numero
    num_bg = add_rounded_rect(slide, x + Inches(0.1), Inches(1.8), Inches(0.5), Inches(0.4), color)
    add_textbox(slide, x + Inches(0.1), Inches(1.82), Inches(0.5), Inches(0.4),
                num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.7), Inches(1.82), Inches(2.1), Inches(0.4),
                phase, font_size=16, color=color, bold=True)

    add_textbox(slide, x + Inches(0.15), Inches(2.35), Inches(2.6), Inches(0.3),
                filename, font_size=11, color=LIGHT_GRAY)

    add_textbox(slide, x + Inches(0.15), Inches(2.7), Inches(2.6), Inches(0.3),
                lines, font_size=11, color=MEDIUM_GRAY)

    txBox = slide.shapes.add_textbox(x + Inches(0.15), Inches(3.1), Inches(2.6), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    for detail in details:
        add_paragraph(tf, f"  > {detail}", 11, WHITE, space_before=Pt(4), space_after=Pt(2))

# Total en bas
total_card = add_rounded_rect(slide, Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.35), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(7.07), Inches(11.7), Inches(0.35),
            "Total : ~950 lignes de code Python  |  Logging complet  |  Gestion d'erreurs  |  Configuration YAML  |  100% reproductible",
            font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 : ANALYSE QUALITE - AVANT NETTOYAGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_RED)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "04", font_size=14, color=ACCENT_RED, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Analyse de Qualite : Etat Initial", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_RED)

# Metriques en haut
metrics = [
    ("50,401", "Lignes", ACCENT_BLUE),
    ("300", "Colonnes", ACCENT_ORANGE),
    ("15.1M", "Cellules totales", LIGHT_GRAY),
    ("26.2%", "Valeurs manquantes", ACCENT_RED),
    ("124 Mo", "Taille fichier", LIGHT_GRAY),
    ("5.5 ans", "Periode (2015-2020)", ACCENT_GREEN),
]

for i, (val, label, color) in enumerate(metrics):
    x = Inches(0.8 + i * 2.05)
    create_metric_card(slide, x, Inches(1.6), Inches(1.85), Inches(1.1), val, label, color)

# Completude
card = add_rounded_rect(slide, Inches(0.8), Inches(3.0), Inches(6.0), Inches(4.2), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(3.1), Inches(5.5), Inches(0.5),
            "Completude par Categorie", font_size=20, color=ACCENT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(3.7), Inches(5.5), Inches(3.3))
tf = txBox.text_frame
tf.word_wrap = True

rows = [
    ("Categorie", "Colonnes", "%", True, ACCENT_BLUE),
    ("Completes (0% manquant)", "92", "30.7%", False, ACCENT_GREEN),
    ("Partielles (<50%)", "189", "63.0%", False, ACCENT_ORANGE),
    ("Majoritairement vides (>=50%)", "19", "6.3%", False, ACCENT_RED),
]

p = tf.paragraphs[0]
p.text = f"{'Categorie':<35} {'Cols':>8} {'%':>8}"
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for cat, cols, pct, is_header, color in rows[1:]:
    add_paragraph(tf, f"{cat:<35} {cols:>8} {pct:>8}", 13, color, space_before=Pt(6))

add_paragraph(tf, "", 8, WHITE, space_before=Pt(12))
add_paragraph(tf, "3,964,527 cellules manquantes sur 15.1M", 14, ACCENT_RED, True)
add_paragraph(tf, "Principalement du aux pays hors focus", 12, LIGHT_GRAY)

# Qualite temporelle
card2 = add_rounded_rect(slide, Inches(7.1), Inches(3.0), Inches(5.5), Inches(4.2), DARK_CARD)
add_textbox(slide, Inches(7.4), Inches(3.1), Inches(5.0), Inches(0.5),
            "Qualite Temporelle", font_size=20, color=ACCENT_GREEN, bold=True)

txBox2 = slide.shapes.add_textbox(Inches(7.4), Inches(3.7), Inches(5.0), Inches(3.3))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "EXCELLENTE"
p.font.size = Pt(20)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.font.name = "Calibri"

checks = [
    ("Gaps temporels", "0", ACCENT_GREEN),
    ("Timestamps dupliques", "0", ACCENT_GREEN),
    ("Resolution", "Horaire (1h)", ACCENT_BLUE),
    ("Regularite", "100%", ACCENT_GREEN),
    ("Debut", "2015-01-01 00:00 UTC", WHITE),
    ("Fin", "2020-06-30 23:00 UTC", WHITE),
]

for label, val, color in checks:
    add_paragraph(tf2, f"  {label}: {val}", 14, color, space_before=Pt(8), space_after=Pt(2))

add_paragraph(tf2, "", 8, WHITE, space_before=Pt(12))
add_paragraph(tf2, "Serie temporelle parfaitement coherente", 14, ACCENT_GREEN, True)
add_paragraph(tf2, "et reguliere sur 5.5 ans", 14, ACCENT_GREEN, True)


# ============================================================
# SLIDE 8 : NETTOYAGE DES DONNEES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "05", font_size=14, color=TEAL, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Nettoyage des Donnees : 4 Etapes", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), TEAL)

# Etape 1
steps = [
    ("1", "Filtrage Pays Focus", ACCENT_BLUE,
     "300 -> 71 colonnes",
     ["Selection DE, DK, FR",
      "Reduction de 76% des colonnes",
      "Focus strategique sur les pays",
      "les plus pertinents"]),
    ("2", "Suppression Colonnes Vides", ACCENT_ORANGE,
     "7 colonnes supprimees",
     ["Seuil : >50% valeurs manquantes",
      "Zone DE_LU (65-67% manquant)",
      "Impact minimal : donnees DE",
      "seules sont completes"]),
    ("3", "Imputation Forward/Backward Fill", ACCENT_GREEN,
     "104,998 valeurs remplies",
     ["Forward fill : tendance recente",
      "Backward fill : debut de serie",
      "Adapte aux series temporelles",
      "Resultat : 0 valeurs manquantes"]),
    ("4", "Variables Temporelles", PURPLE,
     "7 nouvelles variables",
     ["year, month, day, hour",
      "dayofweek, quarter",
      "is_weekend (binaire)",
      "Essentielles pour le ML"]),
]

for i, (num, title, color, subtitle, details) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.7), Inches(2.9), Inches(5.2), DARK_CARD)

    step_bg = add_rounded_rect(slide, x + Inches(0.1), Inches(1.8), Inches(2.7), Inches(0.5), color)
    add_textbox(slide, x + Inches(0.1), Inches(1.83), Inches(2.7), Inches(0.5),
                f"Etape {num}", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.15), Inches(2.5), Inches(2.6), Inches(0.5),
                title, font_size=16, color=color, bold=True)

    add_textbox(slide, x + Inches(0.15), Inches(3.0), Inches(2.6), Inches(0.4),
                subtitle, font_size=12, color=ACCENT_ORANGE, bold=True)

    txBox = slide.shapes.add_textbox(x + Inches(0.15), Inches(3.5), Inches(2.6), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    for detail in details:
        add_paragraph(tf, f"  > {detail}", 12, WHITE, space_before=Pt(5), space_after=Pt(2))


# ============================================================
# SLIDE 9 : RESULTATS AVANT/APRES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "06", font_size=14, color=PURPLE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Resultats : Avant vs Apres Nettoyage", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), PURPLE)

# Tableau comparatif
card = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.4), DARK_CARD)

# En-tete du tableau
header_bg = add_shape(slide, Inches(1.0), Inches(1.9), Inches(11.3), Inches(0.55), ACCENT_BLUE)
headers = [("Metrique", 1.1, 3.0), ("Avant", 4.3, 2.5), ("Apres", 7.0, 2.5), ("Amelioration", 9.5, 2.5)]
for text, x, w in headers:
    add_textbox(slide, Inches(x), Inches(1.95), Inches(w), Inches(0.45),
                text, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Lignes du tableau
rows = [
    ("Colonnes", "300", "71", "-76% (focus strategique)", ACCENT_BLUE),
    ("Valeurs manquantes", "26.2%", "0.0%", "-26.2 points", ACCENT_GREEN),
    ("Taille fichier", "124 Mo", "23 Mo", "-82%", ACCENT_BLUE),
    ("Pays couverts", "32", "3 (DE, DK, FR)", "Focus pertinent", ACCENT_ORANGE),
    ("Lignes conservees", "50,401", "50,401", "100% (aucune perte)", ACCENT_GREEN),
    ("Gaps temporels", "0", "0", "Parfait", ACCENT_GREEN),
    ("Completude", "73.8%", "100%", "+26.2 points", ACCENT_GREEN),
]

for i, (metric, before, after, change, color) in enumerate(rows):
    y = Inches(2.6 + i * 0.6)
    if i % 2 == 0:
        add_shape(slide, Inches(1.0), y, Inches(11.3), Inches(0.55), VERY_DARK)

    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(3.0), Inches(0.45),
                metric, font_size=14, color=WHITE, bold=True)
    add_textbox(slide, Inches(4.3), y + Inches(0.05), Inches(2.5), Inches(0.45),
                before, font_size=14, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.0), y + Inches(0.05), Inches(2.5), Inches(0.45),
                after, font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.5), y + Inches(0.05), Inches(2.5), Inches(0.45),
                change, font_size=13, color=color, alignment=PP_ALIGN.CENTER)

# Message en bas
add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4),
            "Dataset final : 50,401 lignes x 71 colonnes  |  0% valeurs manquantes  |  22.75 Mo  |  Pret pour l'analyse",
            font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10 : PRIX NEGATIFS IDENTIFIES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_YELLOW)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "07", font_size=14, color=ACCENT_YELLOW, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Prix Negatifs Identifies : Variable Cible", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_YELLOW)

# Graphique en barres - Prix negatifs par zone
chart_data = CategoryChartData()
chart_data.categories = ['DK_1', 'DE', 'DK_2', 'FR']
chart_data.add_series('Prix negatifs', (539, 484, 354, 0))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.7),
    Inches(6.0), Inches(3.5), chart_data
)
chart = chart_frame.chart
chart.has_legend = False
plot = chart.plots[0]
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = ACCENT_ORANGE

# Fond du graphique
chart.chart_style = 2
plot.gap_width = 100

# Data labels
series.has_data_labels = True
data_labels = series.data_labels
data_labels.font.size = Pt(14)
data_labels.font.color.rgb = WHITE
data_labels.font.bold = True
data_labels.number_format = '0'
data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END

# Axes
value_axis = chart.value_axis
value_axis.visible = True
value_axis.has_title = False
value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x40, 0x40, 0x60)
tick_labels = value_axis.tick_labels
tick_labels.font.size = Pt(10)
tick_labels.font.color.rgb = LIGHT_GRAY

cat_axis = chart.category_axis
cat_labels = cat_axis.tick_labels
cat_labels.font.size = Pt(12)
cat_labels.font.color.rgb = WHITE
cat_labels.font.bold = True

# Tableau detaille a droite
card = add_rounded_rect(slide, Inches(7.1), Inches(1.7), Inches(5.5), Inches(3.5), DARK_CARD)
add_textbox(slide, Inches(7.4), Inches(1.8), Inches(5.0), Inches(0.5),
            "Detail par Zone de Marche", font_size=18, color=ACCENT_YELLOW, bold=True)

txBox = slide.shapes.add_textbox(Inches(7.4), Inches(2.4), Inches(5.0), Inches(2.6))
tf = txBox.text_frame
tf.word_wrap = True

zones = [
    ("Allemagne (DE)", "484", "2.76%", "-90.01", ACCENT_RED),
    ("Danemark Est (DK_1)", "539", "1.07%", "-58.80", ACCENT_ORANGE),
    ("Danemark Ouest (DK_2)", "354", "0.70%", "-53.62", ACCENT_ORANGE),
    ("France (IT_NORD_FR)", "0", "0.00%", "5.00", ACCENT_GREEN),
]

p = tf.paragraphs[0]
p.text = f"{'Zone':<25} {'Nb':>6} {'%':>8} {'Min':>12}"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for zone, nb, pct, min_val, color in zones:
    add_paragraph(tf, f"{zone:<25} {nb:>6} {pct:>8} {min_val:>8} EUR",
                  12, color, space_before=Pt(8))

# Interpretation en bas
card2 = add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.7), DARK_CARD)
add_textbox(slide, Inches(1.1), Inches(5.6), Inches(5.5), Inches(0.4),
            "Interpretation", font_size=18, color=ACCENT_YELLOW, bold=True)

txBox2 = slide.shapes.add_textbox(Inches(1.1), Inches(6.1), Inches(5.2), Inches(1.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "> 484 exemples en Allemagne = suffisant pour un modele ML de classification binaire robuste"
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.font.name = "Calibri"
add_paragraph(tf2, "> Dataset desepquilibre (~97% positif / ~3% negatif) -> techniques SMOTE, class weights", 13, LIGHT_GRAY)

txBox3 = slide.shapes.add_textbox(Inches(7.0), Inches(6.1), Inches(5.2), Inches(1.0))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "> Allemagne = variable cible principale (plus de donnees)"
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_ORANGE
p.font.bold = True
p.font.name = "Calibri"
add_paragraph(tf3, "> France : 0 prix negatifs -> contraste interessant (nucleaire vs renouvelable)", 13, LIGHT_GRAY)


# ============================================================
# SLIDE 11 : DOCUMENTATION & LIVRABLES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "08", font_size=14, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Documentation & Livrables Produits", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_BLUE)

# Livrable 1 - Dictionnaire
livrables = [
    ("Dictionnaire de Donnees", "docs/dictionnaire_donnees.md", ACCENT_BLUE, "344 lignes",
     ["71 variables documentees", "Types, unites, sources, plages",
      "8 variables temporelles", "3 variables de prix",
      "18 variables de charge", "40+ variables generation",
      "Strategies de nettoyage par variable"]),
    ("Rapport Qualite", "docs/rapport_qualite_donnees.md", ACCENT_GREEN, "465 lignes",
     ["8 sections completes + annexes", "Resume executif",
      "Analyse qualite (avant/apres)", "Prix day-ahead par pays",
      "Strategies de nettoyage", "Limitations et considerations",
      "Recommandations par role"]),
    ("Scripts d'Ingestion", "scripts/01-04_*.py", ACCENT_ORANGE, "~950 lignes",
     ["4 scripts Python modulaires", "Logging complet",
      "Gestion d'erreurs robuste", "Configuration YAML centralisee",
      "100% reproductibles", "Automatisables en pipeline",
      "Code documente et commente"]),
    ("Configuration & Extras", "config/ + reports/", PURPLE, "Fichiers support",
     ["pipeline_config.yaml", "data_quality_report.json",
      "initial_exploration.txt", "requirements.txt",
      "Structure projet organisee", "README.md complet",
      "Fichier .gitignore configure"]),
]

for i, (title, path, color, size, details) in enumerate(livrables):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.6), Inches(2.9), Inches(5.5), DARK_CARD)

    # Badge check
    check_bg = add_rounded_rect(slide, x + Inches(0.1), Inches(1.7), Inches(0.5), Inches(0.4), ACCENT_GREEN)
    add_textbox(slide, x + Inches(0.1), Inches(1.7), Inches(0.5), Inches(0.4),
                "OK", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.7), Inches(1.72), Inches(2.1), Inches(0.4),
                title, font_size=14, color=color, bold=True)

    add_textbox(slide, x + Inches(0.15), Inches(2.2), Inches(2.6), Inches(0.3),
                path, font_size=10, color=LIGHT_GRAY)
    add_textbox(slide, x + Inches(0.15), Inches(2.5), Inches(2.6), Inches(0.3),
                size, font_size=10, color=MEDIUM_GRAY)

    txBox = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.9), Inches(2.6), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    for detail in details:
        add_paragraph(tf, f"  > {detail}", 11, WHITE, space_before=Pt(4), space_after=Pt(2))


# ============================================================
# SLIDE 12 : RECOMMANDATIONS POUR L'EQUIPE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "09", font_size=14, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Recommandations pour l'Equipe", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_GREEN)

# Role 2
card_r2 = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(3.7), Inches(5.2), DARK_CARD)
r2_header = add_rounded_rect(slide, Inches(0.9), Inches(1.8), Inches(3.5), Inches(0.5), ACCENT_GREEN)
add_textbox(slide, Inches(0.9), Inches(1.82), Inches(3.5), Inches(0.5),
            "Role 2 : Analyse Exploratoire", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(3.3), Inches(4.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Charger les donnees :"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

add_paragraph(tf, "pd.read_csv('data/processed/\n  opsd_clean_focus_countries.csv')", 10, LIGHT_GRAY)
add_paragraph(tf, "", 6, WHITE)
add_paragraph(tf, "Visualisations prioritaires :", 12, ACCENT_BLUE, True, space_before=Pt(8))
for item in ["Distribution prix par pays/heure",
             "Correlation generation vs prix",
             "Heatmaps prix negatifs (mois/heure)",
             "Analyse ratio generation/charge",
             "Comparaison profils DE vs DK"]:
    add_paragraph(tf, f"  > {item}", 11, WHITE, space_before=Pt(3))

# Role 3
card_r3 = add_rounded_rect(slide, Inches(4.8), Inches(1.7), Inches(3.7), Inches(5.2), DARK_CARD)
r3_header = add_rounded_rect(slide, Inches(4.9), Inches(1.8), Inches(3.5), Inches(0.5), ACCENT_ORANGE)
add_textbox(slide, Inches(4.9), Inches(1.82), Inches(3.5), Inches(0.5),
            "Role 3 : Feature Engineering", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox2 = slide.shapes.add_textbox(Inches(5.0), Inches(2.5), Inches(3.3), Inches(4.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Features suggerees :"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

features = [
    "renewable_penetration =",
    "  (solar + wind) / load",
    "generation_forecast_error =",
    "  actual - forecast",
    "Moyennes glissantes (3h, 6h, 24h)",
    "Lags temporels (t-1, t-24, t-168)",
    "",
    "Encodage cyclique :",
    "  hour_sin/cos (cycle 24h)",
    "  month_sin/cos (saisonnalite)",
]
for item in features:
    add_paragraph(tf2, f"  {item}", 11, WHITE if ":" not in item else ACCENT_ORANGE,
                  ":" in item and "=" not in item, space_before=Pt(3))

# Role 4
card_r4 = add_rounded_rect(slide, Inches(8.8), Inches(1.7), Inches(3.7), Inches(5.2), DARK_CARD)
r4_header = add_rounded_rect(slide, Inches(8.9), Inches(1.8), Inches(3.5), Inches(0.5), ACCENT_RED)
add_textbox(slide, Inches(8.9), Inches(1.82), Inches(3.5), Inches(0.5),
            "Role 4 : Modelisation", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox3 = slide.shapes.add_textbox(Inches(9.0), Inches(2.5), Inches(3.3), Inches(4.2))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Split temporel recommande :"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

splits = [
    ("Train:", "2015-2018 (70%)", ACCENT_GREEN),
    ("Validation:", "2019 (15%)", ACCENT_ORANGE),
    ("Test:", "2020 (15%)", ACCENT_RED),
]
for label, val, color in splits:
    add_paragraph(tf3, f"  {label} {val}", 12, color, True, space_before=Pt(5))

add_paragraph(tf3, "", 6, WHITE)
add_paragraph(tf3, "Variable cible :", 12, ACCENT_BLUE, True, space_before=Pt(8))
add_paragraph(tf3, "  is_negative_price (binaire)", 11, WHITE, space_before=Pt(3))
add_paragraph(tf3, "  Focus : DE (484 exemples)", 11, ACCENT_ORANGE, space_before=Pt(3))

add_paragraph(tf3, "", 6, WHITE)
add_paragraph(tf3, "Gestion desequilibre :", 12, ACCENT_BLUE, True, space_before=Pt(8))
for item in ["SMOTE / class weights", "Metriques: F1, PR-AUC, ROC-AUC",
             "(pas seulement accuracy)"]:
    add_paragraph(tf3, f"  > {item}", 11, WHITE, space_before=Pt(3))


# ============================================================
# SLIDE 13 : STRUCTURE DU PROJET
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "08b", font_size=14, color=TEAL, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Structure du Projet & Arborescence", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), TEAL)

# Arborescence
card = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(6.5), Inches(5.5), DARK_CARD)
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(6.0), Inches(0.5),
            "Arborescence des Fichiers", font_size=18, color=TEAL, bold=True)

tree_text = """Projet_8_Prix_Negatifs_Electricite/
|-- README.md
|-- LIVRAISON_ROLE_1.md
|-- requirements.txt
|-- config/
|   |-- pipeline_config.yaml
|-- data/
|   |-- raw/opsd_timeseries/
|   |   |-- time_series_60min_singleindex.csv  (124 Mo)
|   |-- processed/
|   |   |-- opsd_clean_focus_countries.csv      (23 Mo)
|   |   |-- opsd_sample_1000.csv
|-- docs/
|   |-- dictionnaire_donnees.md
|   |-- rapport_qualite_donnees.md
|-- scripts/
|   |-- 01_download_opsd_data.py
|   |-- 02_initial_exploration.py
|   |-- 03_data_quality_analysis.py
|   |-- 04_data_cleaning.py
|-- reports/
|   |-- data_quality_report.json
|   |-- initial_exploration.txt
|-- notebooks/                          (pour Role 2)"""

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(6.0), Inches(4.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = tree_text
p.font.size = Pt(10)
p.font.color.rgb = WHITE
p.font.name = "Courier New"

# Technologies utilisees
card2 = add_rounded_rect(slide, Inches(7.6), Inches(1.7), Inches(4.9), Inches(5.5), DARK_CARD)
add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.5),
            "Technologies Utilisees", font_size=18, color=TEAL, bold=True)

techs = [
    ("Python 3", "Langage principal", ACCENT_BLUE),
    ("Pandas", "Manipulation de donnees", ACCENT_GREEN),
    ("NumPy", "Calculs numeriques", ACCENT_GREEN),
    ("Requests", "Telechargement HTTP", ACCENT_ORANGE),
    ("PyYAML", "Configuration pipeline", ACCENT_ORANGE),
    ("Logging", "Tracabilite execution", LIGHT_GRAY),
    ("Git", "Versioning du code", PURPLE),
    ("CSV / JSON", "Formats de donnees", LIGHT_GRAY),
]

txBox2 = slide.shapes.add_textbox(Inches(7.8), Inches(2.4), Inches(4.5), Inches(4.6))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = ""

for tech, desc, color in techs:
    add_paragraph(tf2, f"  {tech}", 14, color, True, space_before=Pt(8), space_after=Pt(1))
    add_paragraph(tf2, f"    {desc}", 11, LIGHT_GRAY, space_before=Pt(1), space_after=Pt(1))


# ============================================================
# SLIDE 14 : TIMELINE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "09b", font_size=14, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Timeline & Chronologie du Travail", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_ORANGE)

weeks = [
    ("S1", "Cadrage", ACCENT_BLUE,
     ["Setup infrastructure projet", "Creation dossiers, config, requirements",
      "Telechargement OPSD Time Series (124 Mo)", "Premiere exploration du dataset"]),
    ("S2", "Exploration", ACCENT_GREEN,
     ["Exploration initiale complete", "Analyse qualite exhaustive",
      "Detection 484 prix negatifs (DE)", "Reunion validation des sources"]),
    ("S3", "Nettoyage", ACCENT_ORANGE,
     ["Nettoyage complet des donnees", "Filtrage 3 pays (300->71 cols)",
      "Forward/backward fill (104,998 val.)", "Creation 7 variables temporelles"]),
    ("S4", "Documentation", ACCENT_RED,
     ["Dictionnaire de donnees (71 vars)", "Rapport qualite (8 sections)",
      "Finalisation des livrables", "Reunion de presentation qualite"]),
]

for i, (week, title, color, tasks) in enumerate(weeks):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.7), Inches(2.9), Inches(4.5), DARK_CARD)

    # Badge semaine
    week_bg = add_rounded_rect(slide, x + Inches(0.1), Inches(1.8), Inches(0.7), Inches(0.5), color)
    add_textbox(slide, x + Inches(0.1), Inches(1.83), Inches(0.7), Inches(0.5),
                week, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.9), Inches(1.83), Inches(1.9), Inches(0.5),
                title, font_size=18, color=color, bold=True)

    # Check mark
    check_bg = add_rounded_rect(slide, x + Inches(2.3), Inches(1.85), Inches(0.4), Inches(0.35), ACCENT_GREEN)
    add_textbox(slide, x + Inches(2.3), Inches(1.85), Inches(0.4), Inches(0.35),
                "OK", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    txBox = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.5), Inches(2.6), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    for task in tasks:
        add_paragraph(tf, f"  > {task}", 11, WHITE, space_before=Pt(6), space_after=Pt(2))

# Barre de progression en bas
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5), DARK_CARD)
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5), ACCENT_GREEN)
add_textbox(slide, Inches(0.8), Inches(6.53), Inches(11.7), Inches(0.5),
            "PHASE ACTIVE TERMINEE (S1-S4)  |  Maintenance S5-S8 : support equipe + ajustements",
            font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15 : COMPETENCES DEMONTREES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "09c", font_size=14, color=PURPLE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Competences Demontrees", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), PURPLE)

competences = [
    ("Data Engineering", ACCENT_BLUE,
     ["Pipeline ETL reproductible", "Gestion donnees volumineuses (124 Mo)",
      "Nettoyage avance series temporelles", "Configuration YAML centralisee"]),
    ("Python & Programmation", ACCENT_GREEN,
     ["950+ lignes de code", "Pandas, NumPy, Requests, PyYAML",
      "Scripts modulaires et reutilisables", "Logging et gestion d'erreurs"]),
    ("Quality Assurance", ACCENT_ORANGE,
     ["Analyse exhaustive de qualite", "Metriques avant/apres nettoyage",
      "Validation temporelle (0 gaps)", "Documentation professionnelle"]),
    ("Communication & Equipe", PURPLE,
     ["Dictionnaire de donnees complet", "Rapport qualite structure",
      "Recommandations actionnables par role", "README et guides d'utilisation"]),
]

for i, (title, color, skills) in enumerate(competences):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.6), Inches(2.9), Inches(3.8), DARK_CARD)

    title_bg = add_rounded_rect(slide, x + Inches(0.1), Inches(1.7), Inches(2.7), Inches(0.5), color)
    add_textbox(slide, x + Inches(0.1), Inches(1.73), Inches(2.7), Inches(0.5),
                title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    txBox = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.4), Inches(2.6), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    for skill in skills:
        add_paragraph(tf, f"  > {skill}", 12, WHITE, space_before=Pt(6), space_after=Pt(2))

# Criteres de validation
card_valid = add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.6), DARK_CARD)
add_textbox(slide, Inches(1.0), Inches(5.7), Inches(3.0), Inches(0.4),
            "Criteres de Validation Remplis", font_size=18, color=ACCENT_GREEN, bold=True)

criteres_s2 = ["Pertinence des donnees confirmee", "Volume suffisant (50k+ lignes)",
               "Granularite adequate (horaire)", "Risques identifies et documentes"]

criteres_s4 = ["Nettoyage termine (0% missing)", "Gestion valeurs manquantes validee",
               "Dictionnaire de donnees complet", "Rapport qualite professionnel"]

txBox_s2 = slide.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(5.5), Inches(1.0))
tf_s2 = txBox_s2.text_frame
tf_s2.word_wrap = True
p = tf_s2.paragraphs[0]
p.text = "Semaine 2 - Validation Sources"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"
for c in criteres_s2:
    add_paragraph(tf_s2, f"  [OK] {c}", 10, ACCENT_GREEN, space_before=Pt(2))

txBox_s4 = slide.shapes.add_textbox(Inches(7.0), Inches(6.1), Inches(5.5), Inches(1.0))
tf_s4 = txBox_s4.text_frame
tf_s4.word_wrap = True
p = tf_s4.paragraphs[0]
p.text = "Semaine 4 - Qualite des Donnees"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_ORANGE
p.font.bold = True
p.font.name = "Calibri"
for c in criteres_s4:
    add_paragraph(tf_s4, f"  [OK] {c}", 10, ACCENT_GREEN, space_before=Pt(2))


# ============================================================
# SLIDE 16 : CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(1), Inches(0.5),
            "10", font_size=14, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
            "Conclusion & Prochaines Etapes", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_GREEN)

# Points forts
card_pf = add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(3.0), DARK_CARD)
add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5.4), Inches(0.5),
            "Points Forts du Travail", font_size=20, color=ACCENT_GREEN, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.4), Inches(2.1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ""
for item in [
    "Source academique de haute qualite (TU Berlin / ETH Zurich)",
    "Coherence temporelle parfaite : 0 gaps, 0 doublons",
    "Completude finale 100% : 0 valeurs manquantes",
    "484 prix negatifs en Allemagne : phenomene confirme",
    "Documentation exhaustive pour toute l'equipe",
]:
    add_paragraph(tf, f"  > {item}", 13, WHITE, space_before=Pt(5), space_after=Pt(2))

# Livrables recapitulatif
card_liv = add_rounded_rect(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(3.0), DARK_CARD)
add_textbox(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
            "Recapitulatif des Livrables", font_size=20, color=ACCENT_BLUE, bold=True)

livrables_recap = [
    ("Donnees nettoyees", "23 Mo, 0% missing", ACCENT_GREEN),
    ("4 scripts Python", "~950 lignes, reproductibles", ACCENT_ORANGE),
    ("Dictionnaire de donnees", "71 variables documentees", ACCENT_BLUE),
    ("Rapport qualite", "8 sections completes", ACCENT_BLUE),
    ("Configuration pipeline", "YAML centralise", LIGHT_GRAY),
]

txBox2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(2.1))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = ""
for item, detail, color in livrables_recap:
    add_paragraph(tf2, f"  [OK] {item}", 13, color, True, space_before=Pt(5))
    add_paragraph(tf2, f"        {detail}", 11, LIGHT_GRAY, space_before=Pt(1))

# Prochaines etapes
card_next = add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.4), DARK_CARD)
add_textbox(slide, Inches(1.0), Inches(4.9), Inches(5.0), Inches(0.5),
            "Prochaines Etapes (S5-S8)", font_size=20, color=ACCENT_ORANGE, bold=True)

txBox3 = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(5.5), Inches(1.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Maintenance & Support :"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for item in ["Support Role 2 pour questions donnees",
             "Ajustements si problemes identifies en EDA",
             "Documentation continue"]:
    add_paragraph(tf3, f"  > {item}", 12, WHITE, space_before=Pt(4))

txBox4 = slide.shapes.add_textbox(Inches(7.0), Inches(5.5), Inches(5.2), Inches(1.5))
tf4 = txBox4.text_frame
tf4.word_wrap = True
p = tf4.paragraphs[0]
p.text = "Ameliorations possibles :"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"

for item in ["Telecharger OPSD Weather Data (ERA5)",
             "Obtenir ENTSO-E API token si besoin",
             "Dashboard interactif qualite des donnees"]:
    add_paragraph(tf4, f"  > {item}", 12, WHITE, space_before=Pt(4))


# ============================================================
# SLIDE 17 : MERCI / QUESTIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
            "Merci !", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5), Inches(3.2), Inches(3.333), Inches(0.04), ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
            "Questions ?", font_size=28, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

# Metriques recap
metrics_final = [
    ("50,401", "Observations", ACCENT_BLUE),
    ("71", "Variables", ACCENT_GREEN),
    ("0%", "Missing", ACCENT_GREEN),
    ("484", "Prix negatifs", ACCENT_ORANGE),
    ("~950", "Lignes Python", PURPLE),
]

for i, (val, label, color) in enumerate(metrics_final):
    x = Inches(1.5 + i * 2.2)
    create_metric_card(slide, x, Inches(4.6), Inches(1.8), Inches(1.1), val, label, color)

add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
            "Projet 8 : Prix Negatifs de l'Electricite Renouvelable  |  Role 1 : Data Engineering  |  Fevrier 2026",
            font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAUVEGARDE
# ============================================================
output_path = "/Users/loulou/Documents/Documents - Mac'Donald/school/s6/projet 2eme session/sujet-support/.claude/worktrees/ecstatic-perlman/Projet_8_Prix_Negatifs_Electricite/Presentation_Role_1_Data_Engineering.pptx"
prs.save(output_path)
print(f"Presentation sauvegardee : {output_path}")
print(f"Nombre de slides : {len(prs.slides)}")
